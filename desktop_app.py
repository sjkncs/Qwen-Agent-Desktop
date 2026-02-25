#!/usr/bin/env python3
"""
Qwen-Agent Desktop — Native-quality desktop application
Architecture: aiohttp local server + SPA frontend + SSE streaming

Usage:
    python desktop_app.py
    python desktop_app.py --model gemini-2.5-pro
    python desktop_app.py --port 9720
"""

import argparse
import asyncio
import json
import os
import sys
import threading
import webbrowser

from aiohttp import web

# ── Import the data/conversation manager ──
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from qwen_agent.gui.desktop.api_bridge import ApiBridge

# Support both dev mode and PyInstaller bundle
_BASE = getattr(sys, '_MEIPASS', os.path.dirname(os.path.abspath(__file__)))
FRONTEND_DIR = os.path.join(_BASE, 'qwen_agent', 'gui', 'desktop')


def create_app(bridge: ApiBridge) -> web.Application:
    """Build the aiohttp application with API routes + static serving."""

    app = web.Application()

    # ── API Routes ─────────────────────────────────────────────

    async def api_conversations(request):
        return web.json_response(json.loads(bridge.get_conversations()))

    async def api_new_conversation(request):
        return web.json_response(json.loads(bridge.new_conversation()))

    async def api_switch_conversation(request):
        data = await request.json()
        msgs = json.loads(bridge.switch_conversation(data['id']))
        return web.json_response(msgs)

    async def api_delete_conversation(request):
        data = await request.json()
        bridge.delete_conversation(data['id'])
        return web.json_response({'ok': True})

    async def api_models(request):
        return web.json_response(json.loads(bridge.get_models()))

    async def api_set_model(request):
        data = await request.json()
        bridge.set_model(data['model'])
        return web.json_response({'ok': True})

    async def api_current_model(request):
        return web.json_response(json.loads(bridge.get_current_model()))

    async def api_system_info(request):
        return web.json_response(json.loads(bridge.get_system_info()))

    async def api_current_conv_id(request):
        return web.json_response(json.loads(bridge.get_current_conv_id()))

    async def api_chat_stream(request):
        """SSE endpoint: streams tokens as server-sent events."""
        data = await request.json()
        text = data.get('text', '').strip()
        mode = data.get('mode', 'chat')
        model = data.get('model')
        ephemeral = data.get('ephemeral', False)

        if not text:
            return web.json_response({'error': 'empty'}, status=400)

        use_model = model or bridge._default_model
        system_prompt = bridge._get_system_prompt(mode)

        if ephemeral:
            # Ephemeral mode: don't touch conversation history
            api_messages = []
            if system_prompt:
                api_messages.append({'role': 'system', 'content': system_prompt})
            api_messages.append({'role': 'user', 'content': text})
        else:
            # Normal mode: save to conversation
            if not bridge._current_conv_id:
                bridge.new_conversation()

            conv = bridge._conversations[bridge._current_conv_id]
            conv['messages'].append({'role': 'user', 'content': text})

            # Auto-title
            if sum(1 for m in conv['messages'] if m['role'] == 'user') == 1:
                conv['title'] = text[:30] + ('…' if len(text) > 30 else '')

            use_model = model or conv.get('model', bridge._default_model)
            bridge._save_conversations()

            api_messages = []
            if system_prompt:
                api_messages.append({'role': 'system', 'content': system_prompt})
            api_messages.extend([
                {'role': m['role'], 'content': m['content']}
                for m in conv['messages']
            ])

        # Set up SSE response
        response = web.StreamResponse(
            status=200,
            reason='OK',
            headers={
                'Content-Type': 'text/event-stream',
                'Cache-Control': 'no-cache',
                'Connection': 'keep-alive',
                'X-Accel-Buffering': 'no',
            },
        )
        await response.prepare(request)

        # Stream from OpenAI-compatible API
        full_response = ''
        try:
            import openai
            client = openai.OpenAI(api_key=bridge._api_key, base_url=bridge._api_base)

            # Title event (only for persistent chats)
            if not ephemeral and bridge._current_conv_id:
                conv = bridge._conversations[bridge._current_conv_id]
                await response.write(
                    f"event: title\ndata: {json.dumps({'id': bridge._current_conv_id, 'title': conv['title']})}\n\n".encode()
                )

            stream = client.chat.completions.create(
                model=use_model,
                messages=api_messages,
                stream=True,
                max_tokens=4096,
            )

            for chunk in stream:
                if not chunk.choices:
                    continue
                delta = chunk.choices[0].delta
                if delta and delta.content:
                    full_response += delta.content
                    await response.write(
                        f"event: token\ndata: {json.dumps(delta.content)}\n\n".encode()
                    )

            # Save assistant response (only for persistent chats)
            if not ephemeral and bridge._current_conv_id:
                conv = bridge._conversations[bridge._current_conv_id]
                conv['messages'].append({'role': 'assistant', 'content': full_response})
                bridge._save_conversations()

        except Exception as e:
            await response.write(
                f"event: error\ndata: {json.dumps(str(e))}\n\n".encode()
            )

        await response.write(b"event: done\ndata: {}\n\n")
        await response.write_eof()
        return response

    # ── File Upload + Parse ───────────────────────────────────

    async def api_upload(request):
        """Parse uploaded file (PDF/DOCX/PPTX/image/text) and return extracted text."""
        import tempfile
        reader = await request.multipart()
        field = await reader.next()
        if not field:
            return web.json_response({'error': 'No file'}, status=400)

        filename = field.filename or 'unknown'
        data = await field.read(decode=False)
        ext = os.path.splitext(filename)[1].lower()

        text = ''
        file_type = 'text'
        meta = {'name': filename, 'size': len(data)}

        try:
            if ext == '.pdf':
                file_type = 'pdf'
                import fitz  # PyMuPDF
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                    tmp.write(data)
                    tmp_path = tmp.name
                try:
                    doc = fitz.open(tmp_path)
                    meta['pages'] = len(doc)
                    pages_text = []
                    for i, page in enumerate(doc):
                        pt = page.get_text('text')
                        if pt.strip():
                            pages_text.append(f'--- 第 {i+1} 页 ---\n{pt.strip()}')
                    text = '\n\n'.join(pages_text)
                    doc.close()
                finally:
                    os.unlink(tmp_path)

            elif ext in ('.docx',):
                file_type = 'docx'
                import docx
                import io
                doc = docx.Document(io.BytesIO(data))
                meta['paragraphs'] = len(doc.paragraphs)
                parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
                # Also extract tables
                for table in doc.tables:
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        parts.append(' | '.join(cells))
                text = '\n'.join(parts)

            elif ext in ('.pptx',):
                file_type = 'pptx'
                from pptx import Presentation
                import io
                prs = Presentation(io.BytesIO(data))
                meta['slides'] = len(prs.slides)
                parts = []
                for i, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    slide_texts.append(para.text.strip())
                    if slide_texts:
                        parts.append(f'--- 幻灯片 {i+1} ---\n' + '\n'.join(slide_texts))
                text = '\n\n'.join(parts)

            elif ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.svg', '.tiff', '.ico'):
                file_type = 'image'
                from PIL import Image
                import io, base64
                try:
                    img = Image.open(io.BytesIO(data))
                    meta['width'] = img.width
                    meta['height'] = img.height
                    meta['format'] = img.format
                    # Convert to base64 for sending to vision-capable models
                    buf = io.BytesIO()
                    if img.mode in ('RGBA', 'P'):
                        img = img.convert('RGB')
                    img.save(buf, format='JPEG', quality=85)
                    b64 = base64.b64encode(buf.getvalue()).decode()
                    meta['base64'] = b64
                    text = f'[图片: {filename}, {img.width}x{img.height}px, {img.format}]'
                except Exception:
                    text = f'[无法解析图片: {filename}]'

            elif ext in ('.xls', '.xlsx', '.csv'):
                file_type = 'spreadsheet'
                if ext == '.csv':
                    text = data.decode('utf-8', errors='replace')
                else:
                    text = f'[Excel文件: {filename}, 请使用CSV格式上传以获得最佳解析效果]'

            else:
                # Plain text files
                file_type = 'text'
                text = data.decode('utf-8', errors='replace')

        except Exception as e:
            text = f'[文件解析错误: {str(e)}]'

        if len(text) > 15000:
            text = text[:15000] + '\n\n… (内容已截断，共约' + str(len(text)) + '字符)'

        return web.json_response({
            'text': text,
            'type': file_type,
            'meta': meta,
        })

    # ── Static file serving ────────────────────────────────────

    async def index_handler(request):
        return web.FileResponse(os.path.join(FRONTEND_DIR, 'index.html'))

    # ── Register routes ────────────────────────────────────────

    async def favicon_handler(request):
        icon = os.path.join(FRONTEND_DIR, 'Qwen3.png')
        if os.path.exists(icon):
            return web.FileResponse(icon, headers={'Content-Type': 'image/png'})
        return web.Response(status=204)

    app.router.add_get('/', index_handler)
    app.router.add_get('/favicon.ico', favicon_handler)
    app.router.add_get('/api/conversations', api_conversations)
    app.router.add_post('/api/conversations/new', api_new_conversation)
    app.router.add_post('/api/conversations/switch', api_switch_conversation)
    app.router.add_post('/api/conversations/delete', api_delete_conversation)
    app.router.add_get('/api/models', api_models)
    app.router.add_post('/api/model', api_set_model)
    app.router.add_get('/api/model/current', api_current_model)
    app.router.add_get('/api/system-info', api_system_info)
    app.router.add_get('/api/current-conv-id', api_current_conv_id)
    app.router.add_post('/api/chat', api_chat_stream)
    app.router.add_post('/api/upload', api_upload)
    app.router.add_static('/', FRONTEND_DIR, show_index=False)

    return app


def main():
    parser = argparse.ArgumentParser(description='Qwen-Agent Desktop')
    parser.add_argument('--model', default=os.environ.get('QWEN_MODEL', 'claude-sonnet-4-5'), help='Default model')
    parser.add_argument('--api-base', default=os.environ.get('QWEN_API_BASE', 'https://hiapi.online/v1'))
    parser.add_argument('--api-key', default=os.environ.get('QWEN_API_KEY', ''))
    parser.add_argument('--port', type=int, default=9720, help='Local server port')
    parser.add_argument('--no-browser', action='store_true', help='Don\'t open browser')
    args = parser.parse_args()

    bridge = ApiBridge(
        api_key=args.api_key,
        api_base=args.api_base,
        default_model=args.model,
    )

    app = create_app(bridge)
    url = f'http://localhost:{args.port}'

    print('=' * 50)
    print('  Qwen-Agent Desktop')
    print('=' * 50)
    print(f'  模型    : {args.model}')
    print(f'  API     : {args.api_base}')
    print(f'  地址    : {url}')
    print(f'  快捷键  : Ctrl+N 新对话 | Ctrl+B 侧边栏 | Ctrl+D 主题')
    print('=' * 50)

    if not args.no_browser:
        threading.Timer(1.5, lambda: webbrowser.open(url)).start()

    web.run_app(app, host='127.0.0.1', port=args.port, print=None)


if __name__ == '__main__':
    main()
